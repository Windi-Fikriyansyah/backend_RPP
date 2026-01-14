import httpx
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy
import os

class PPTService:
    TEMPLATE_DIR = "app/templates"

    @staticmethod
    def _replace_text_in_shape(shape, replacements):
        if not shape.has_text_frame:
            return
            
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                # Create a local copy to modify
                new_text = original_text
                for key, val in replacements.items():
                    if key in new_text:
                        new_text = new_text.replace(key, str(val))
                
                # Only update if changed to preserve formatting where possible
                if new_text != original_text:
                    run.text = new_text

    @staticmethod
    def _duplicate_slide(prs, source_slide_index):
        """
        Duplicate a slide by copying the layout and all shapes.
        This is a workaround since python-pptx doesn't support cloning directly.
        """
        source_slide = prs.slides[source_slide_index]
        layout = source_slide.slide_layout
        new_slide = prs.slides.add_slide(layout)
        
        for shape in source_slide.shapes:
            # We only copy TextBoxes and Pictures for now to keep it simple and fast
            # Copying complex shapes is very hard in python-pptx without deep XML manipulation
            
            if shape.has_text_frame:
                # Add new textbox with same position
                new_shape = new_slide.shapes.add_textbox(
                    shape.left, shape.top, shape.width, shape.height
                )
                
                # Copy text/formatting
                # Simple copy:
                new_shape.text_frame.text = shape.text_frame.text
                
                # Try to copy style of first paragraph
                if shape.text_frame.paragraphs:
                    p_src = shape.text_frame.paragraphs[0]
                    p_dst = new_shape.text_frame.paragraphs[0]
                    if p_src.runs:
                        r_src = p_src.runs[0]
                        if p_dst.runs:
                            r_dst = p_dst.runs[0]
                            r_dst.font.size = r_src.font.size
                            r_dst.font.bold = r_src.font.bold
                            r_dst.font.color.rgb = r_src.font.color.rgb
                            try:
                                r_dst.font.name = r_src.font.name
                            except: pass
            
            # Note: We skip images in the template for duplication to avoid bloating 
            # unless they are background shapes (which are usually on the master).
            # If the user wants a logo on every slide, it should be on the Master Slide.

        return new_slide
    
    @classmethod
    async def generate_ppt(cls, json_data: dict) -> BytesIO:
        theme_name = json_data.get("theme", "Ceria")
        template_filename = f"{theme_name}.pptx"
        template_path = os.path.join(cls.TEMPLATE_DIR, template_filename)
        
        # Fallback to Ceria if specific theme not found
        if not os.path.exists(template_path):
            print(f"Template {theme_name} not found, using Ceria")
            template_path = os.path.join(cls.TEMPLATE_DIR, "Ceria.pptx")
            
        print(f"DEBUG: Loading template from {template_path}")
        prs = Presentation(template_path)
        
        # --- SLIDE 1: TITLE (Index 0) ---
        if len(prs.slides) > 0:
            title_slide = prs.slides[0]
            shape_map_title = {}
            cls._replace_text_in_shape_recursive(title_slide, {
                "{{judul_materi}}": json_data.get("judul_materi", ""),
                "{{theme}}": json_data.get("theme", "")
            }, shape_map=shape_map_title)
            
            # Dynamic Font Sizing for Main Title (User Request)
            try:
                if "{{judul_materi}}" in shape_map_title:
                    main_title_shape = shape_map_title["{{judul_materi}}"]
                    if main_title_shape.has_text_frame:
                        main_title_shape.text_frame.word_wrap = True
                        raw_main_title = json_data.get("judul_materi", "")
                        words_main = raw_main_title.split()
                        
                        if len(words_main) > 6:
                             # Shrink Main Title
                             new_size_main = Pt(36) # Default is usually 44-54
                             for p in main_title_shape.text_frame.paragraphs:
                                for r in p.runs:
                                    r.font.size = new_size_main
                                    r.font.bold = True # Enforce Bold
            except Exception as e:
                print(f"DEBUG: Main Title Font adjustment failed: {e}")
        
        # --- CONTENT SLIDES (Index 1..N) ---
        slides_data = json_data.get("slides", [])
        
        # Max slides we can fill is determined by the template (minus title slide)
        available_slots = len(prs.slides) - 1
        
        if available_slots < 1:
            # Fallback if template is broken/empty
            print("Template has no content slides!")
        else:
            # Fill existing slides
            # We strictly map Item 0 -> Slide 1, Item 1 -> Slide 2...
            # We stop if we run out of items OR run out of slides.
            
            items_to_process = slides_data[:available_slots] # Truncate excess items
            
            for idx, slide_data in enumerate(items_to_process):
                # Target slide index is idx + 1 (because 0 is title)
                target_slide = prs.slides[idx + 1]
                
                shape_map = {}
                slide_content = "\n".join([f"• {x}" for x in slide_data.get("konten", [])])
                
                cls._replace_text_in_shape_recursive(target_slide, {
                    "{{judul_slide}}": slide_data.get("judul_slide", ""),
                    "{{konten}}": slide_content
                }, shape_map=shape_map)
                
                # Dynamic Font Sizing for Title (User Request)
                try:
                    if "{{judul_slide}}" in shape_map:
                        title_shape = shape_map["{{judul_slide}}"]
                        if title_shape.has_text_frame:
                            title_shape.text_frame.word_wrap = True
                            
                            # Check source data directly for reliability
                            raw_title = slide_data.get("judul_slide", "")
                            words = raw_title.split()
                            word_count = len(words)
                            
                            new_size = None
                            if word_count > 12:
                                new_size = Pt(20) # Very long title
                            elif word_count > 4: # Changed from 5 to 4 per request
                                new_size = Pt(24) # Medium-long title
                                
                            if new_size:
                                for p in title_shape.text_frame.paragraphs:
                                    for r in p.runs:
                                        r.font.size = new_size
                                        r.font.bold = True # Enforce Bold
                                        
                    # Enforce Regular Font for Content (User Request)
                    if "{{konten}}" in shape_map:
                        content_shape = shape_map["{{konten}}"]
                        if content_shape.has_text_frame:
                             for p in content_shape.text_frame.paragraphs:
                                for r in p.runs:
                                    r.font.bold = False # Enforce Regular
                                    
                except Exception as e:
                    print(f"DEBUG: Font adjustment failed: {e}")
                except Exception as e:
                    print(f"DEBUG: Font adjustment failed: {e}")

            # --- DELETE UNUSED SLIDES ---
            # If we used X items, we used indices 1 to X.
            # Unused indices are X+1 to End.
            # We must delete them BACKWARDS to avoid index shifting.
            
            used_count = len(items_to_process)
            total_slides = len(prs.slides) # This includes title (Index 0)
            
            # Last used index was `used_count` (0 is title, 1..used_count are content).
            # So first unused index is `used_count + 1`.
            
            first_unused_index = used_count + 1
            
            if first_unused_index < total_slides:
                # Iterate backwards from last index down to first_unused_index
                print(f"DEBUG: Deleting unused slides from {first_unused_index} to {total_slides-1}")
                
                # Get XML Key for deletion
                xml_slides = prs.slides._sldIdLst
                
                slides_to_delete = []
                for i in range(total_slides - 1, first_unused_index - 1, -1):
                    slides_to_delete.append(prs.slides[i].slide_id)

                # Delete logic
                for slide_id in slides_to_delete:
                    # Find and remove
                    for i, sldId in enumerate(xml_slides):
                        if sldId.id == slide_id:
                            xml_slides.remove(sldId)
                            break

        ppt_output = BytesIO()
        prs.save(ppt_output)
        ppt_output.seek(0)
        return ppt_output

    @staticmethod
    def _replace_text_in_shape_recursive(slide_or_group, replacements, shape_map=None):
        # Handle both Slide and Group objects which have .shapes to iterate
        shapes = slide_or_group.shapes
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        for shape in shapes:
            # 1. Check if it's a Group -> Recurse
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                PPTService._replace_text_in_shape_recursive(shape, replacements, shape_map)
                continue
            
            # 2. Check if it has text frame
            if shape.has_text_frame:
                # Iterate paragraphs
                for p in shape.text_frame.paragraphs:
                    full_text = "".join(r.text for r in p.runs)
                    if not full_text.strip():
                        continue

                    original = full_text
                    modified = False
                    
                    for k, v in replacements.items():
                        if k in full_text:
                            full_text = full_text.replace(k, str(v))
                            modified = True
                            # Track matched shape
                            if shape_map is not None:
                                shape_map[k] = shape
                            
                    if modified:
                        # Capture style from the FIRST run
                        first_run_font_size = None
                        first_run_font_color = None
                        first_run_bold = None
                        
                        if p.runs:
                            r0 = p.runs[0]
                            first_run_font_size = r0.font.size
                            first_run_bold = r0.font.bold
                            try:
                                first_run_color = r0.font.color.rgb
                            except:
                                first_run_color = None

                        # CRITICAL FIX: Clear ALL runs in the paragraph to prevent overlapping/appending
                        p.clear() 
                        
                        # Add new run with the FULL replaced text
                        new_run = p.add_run()
                        new_run.text = full_text
                        
                        # Restore style
                        if first_run_font_size:
                            new_run.font.size = first_run_font_size
                        if first_run_font_color:
                            new_run.font.color.rgb = first_run_color
                        if first_run_bold is not None:
                            new_run.font.bold = first_run_bold

    @staticmethod
    def _duplicate_slide_native(prs, source_slide):
        """
        Duplicate a slide by creating a new slide with the same layout
        and then copying all shape elements via XML deep copy.
        CRITICAL: Also copies image relationships (blips).
        """
        layout = source_slide.slide_layout
        new_slide = prs.slides.add_slide(layout)
        
        # CLEAR DEFAULT PLACEHOLDERS
        # When we add a slide with a layout, it might come with empty placeholders.
        # Since we are about to copy all shapes from the source (which includes the filled placeholders),
        # we should remove the default empty ones to avoid "Text Overlap" (Double Textboxes).
        for shp in list(new_slide.shapes):
            sp = shp.element
            sp.getparent().remove(sp)

        # Copy every shape from source to destination via XML
        for shape in source_slide.shapes:
            new_element = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
            
        # FIX MISSING IMAGES (Relatioships)
        source_rels = source_slide.part.rels
        dest_part = new_slide.part
        
        blip_list = new_slide.shapes._spTree.xpath('.//a:blip')
        for blip in blip_list:
            rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if rId and rId in source_rels:
                rel = source_rels[rId]
                if "image" in rel.reltype:
                    new_rId = dest_part.relate_to(rel.target_part, rel.reltype)
                    blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)

        return new_slide

